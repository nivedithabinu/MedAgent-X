import os
import json
import io
import time
import uuid
import logging
import traceback
import numpy as np
from typing import List, Dict, Any, Optional

from google import genai
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv
from pypdf import PdfReader
from fastapi.responses import StreamingResponse
from pptx import Presentation

load_dotenv(override=True)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

app = FastAPI(title="MedAgent-X API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

class QueryRequest(BaseModel):
    query: str
    doc_id: str

class DocRequest(BaseModel):
    doc_id: str

class VectorDatabase:
    """In-memory Vector Store mimicking enterprise solutions like ChromaDB/Pinecone."""
    def __init__(self):
        self._store: Dict[str, Dict[str, Any]] = {}

    def save_document(self, doc_id: str, chunks: List[str], embeddings: List[np.ndarray], full_text: str, filename: str, pages: int):
        self._store[doc_id] = {
            "chunks": chunks,
            "embeddings": embeddings,
            "full_text": full_text,
            "filename": filename,
            "pages": pages,
            "ppt": None,
            "graph": None
        }

    def get_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        return self._store.get(doc_id)

    def update_document(self, doc_id: str, key: str, value: Any):
        if doc_id in self._store:
            self._store[doc_id][key] = value

# Singleton DB instance
db = VectorDatabase()

class AIService:
    def __init__(self):
        raw_key = os.getenv("GEMINI_API_KEY", "")
        self.api_key = raw_key.replace('"', '').replace("'", "").replace("Bearer ", "").strip()
        self.demo_mode = not bool(self.api_key)
        
        if not self.demo_mode:
            try:
                self.client = genai.Client(api_key=self.api_key)
                logger.info("AI Service initialized with Gemini API.")
            except Exception as e:
                logger.error(f"Failed to initialize GenAI client: {e}. Defaulting to Demo Mode.")
                self.demo_mode = True
        else:
            logger.warning("No valid GEMINI_API_KEY found. AI Service starting in DEMO MODE.")

    def generate_content(self, prompt: str, instruction: str, is_json: bool = False) -> str:
        if self.demo_mode:
            time.sleep(1.5)
            if is_json:
                return """[
                    {"title": "Summary", "bullets": ["Clinical trial for novel Heart Failure (HF) therapeutic.", "Study involved 1,200 participants over 24 months.", "Primary endpoint: Reduction in hospitalizations."], "icon": "ph-heartbeat"},
                    {"title": "Methodology", "bullets": ["Double-blind, randomized placebo-controlled trial.", "Patients administered 10mg daily dose.", "Monitored via regular ECG and blood biomarkers."], "icon": "ph-flask"},
                    {"title": "Key Findings", "bullets": ["34% reduction in cardiovascular events.", "Significant improvement in Left Ventricular Ejection Fraction (LVEF).", "P-value < 0.001 demonstrating high statistical significance."], "icon": "ph-chart-line-up"},
                    {"title": "Safety & Tolerability", "bullets": ["Adverse events were comparable to the placebo group.", "No instances of severe hepatotoxicity observed.", "Mild nausea reported in 4% of participants."], "icon": "ph-shield-check"},
                    {"title": "Conclusion", "bullets": ["The therapeutic presents a breakthrough in HF management.", "FDA fast-track approval recommended based on efficacy.", "Phase IV post-marketing surveillance to follow."], "icon": "ph-check-circle"}
                ]"""

            elif "mindmap" in prompt.lower():
                return """mindmap
  root((Heart Failure))
    Symptoms
      Dyspnea
      Fatigue
      Edema
    Diagnostics
      Echocardiogram
      BNP Blood Test
      Chest X-Ray
    Treatments
      ACE Inhibitors
      Beta Blockers
      Diuretics"""

            elif "strict medical classifier" in instruction.lower():
                return "YES"

            else:
                return "Based on the clinical documentation, the treatment showed a statistically significant improvement in patient outcomes (p < 0.001), specifically reducing hospital readmission rates by 34% over a 24-month period. [PAGE 4]"

        # --- REAL API CALL ---
        mime_type = "application/json" if is_json else "text/plain"
        for attempt in range(3):
            try:
                response = self.client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=prompt,
                    config=genai.types.GenerateContentConfig(
                        system_instruction=instruction,
                        temperature=0.2,
                        response_mime_type=mime_type
                    )
                )

                return response.text

            except Exception as e:
                logger.warning(f"Generation attempt {attempt + 1} failed: {e}")

                if "401" in str(e) or "403" in str(e) or "429" in str(e):
                    logger.error("Authentication or Quota failed. Falling back to Demo Mode.")
                    self.demo_mode = True
                    return self.generate_content(prompt, instruction, is_json)

                if attempt == 2:
                    raise Exception(f"Failed to generate content: {str(e)}")

                time.sleep(2 * (attempt + 1))

    def generate_embedding(self, text: str) -> np.ndarray:
        if self.demo_mode:
            return np.random.rand(768)

        for attempt in range(3):
            try:
                response = self.client.models.embed_content(
                    model="gemini-embedding-2",
                    contents=text.strip()
                )

                return np.array(response.embeddings[0].values)

            except Exception as e:
                if "429" in str(e):
                    time.sleep(2)
                    continue

                if "401" in str(e) or "403" in str(e):
                    self.demo_mode = True
                    return self.generate_embedding(text)

                if attempt == 2:
                    raise Exception(f"Embedding failed: {str(e)}")

ai_service = AIService()

@app.get("/")
def read_root():
    return {
        "status": "MedAgent-X Enterprise Backend Online",
        "demo_mode_active": ai_service.demo_mode,
        "version": "2.0.0 (Enterprise SOA)"
    }

@app.post("/api/upload")
async def upload_pdf(file: UploadFile = File(...)):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Must be a PDF file.")
        
    try:
        content = await file.read()
        pdf = PdfReader(io.BytesIO(content))
        full_text = ""

        # Defensive PDF Extraction
        for i, page in enumerate(pdf.pages): 
            try:
                text = page.extract_text()
                if text:
                    full_text += f"\n--- [PAGE {i+1}] ---\n{text.strip()}\n"
    
            except Exception as e:
                logger.warning(f"Skipping page {i+1} due to extraction error: {e}")
                continue

        if not full_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from PDF.")

        # Agentic Verification Step
        first_page = full_text[:1000]
        prompt = f"Analyze this text snippet: \"{first_page}\". Is this highly likely related to Medical, Health, or Biological sciences? Reply with exactly YES or NO."
        
        try:
            response_text = ai_service.generate_content(prompt, "You are a strict medical classifier.")

        except Exception as e:
            logger.warning(f"Classification failed ({e}), proceeding anyway.")
            response_text = "YES"
        
        if "YES" not in response_text.strip().upper():
            raise HTTPException(status_code=403, detail="Document rejected: Content does not appear to be medical research.")

        # Overlapping Semantic Chunking (Advanced RAG Technique)
        chunk_size = 1500
        overlap = 200
        chunks = []
        start = 0

        while start < len(full_text):
            end = start + chunk_size
            chunks.append(full_text[start:end].strip())
            start += chunk_size - overlap
        
        chunks = [c for c in chunks if c][:20] # Limit to 20 for processing speed

        try:
            embeddings = []
            for i, chunk in enumerate(chunks):
                emb = ai_service.generate_embedding(chunk)
                embeddings.append(emb)
                time.sleep(0.3) # Rate limit safety

        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

        doc_id = str(uuid.uuid4())
        db.save_document(doc_id, chunks, embeddings, full_text[:25000], file.filename, len(pdf.pages))

        return {
            "doc_id": doc_id, 
            "filename": file.filename, 
            "pages_count": len(pdf.pages)
        }

    except HTTPException as e:
        raise e

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.post("/api/chat")
async def chat_with_agent(req: QueryRequest):
    doc_data = db.get_document(req.doc_id)
    if not doc_data:
        raise HTTPException(status_code=404, detail="Session expired. Please re-upload the document.")
    
    try:
        # Custom Vector Search (Cosine Similarity)
        query_emb = ai_service.generate_embedding(req.query)
        query_emb_np = np.array(query_emb)
        doc_embs_np = np.array(doc_data["embeddings"])
        
        similarities = [np.dot(query_emb_np, doc_emb) / (np.linalg.norm(query_emb_np) * np.linalg.norm(doc_emb) + 1e-10) for doc_emb in doc_embs_np]
        top_indices = np.argsort(similarities)[-4:][::-1]
        relevant_context = "\n\n...\n\n".join([doc_data["chunks"][i] for i in top_indices])
        
        prompt = f"User Query: {req.query}\n\nRelevant Document Context:\n{relevant_context}"
        instruction = """You are MedAgent-X, an advanced Clinical AI assistant.
        1. If the query relates to the provided Document Context, answer using the context and explicitly cite the [PAGE X] markers.
        2. If the query is a general question (e.g., coding, math, casual conversation), ignore the document and answer using your broad general knowledge.
        3. Always be professional, clear, and highly helpful."""
        
        reply_text = ai_service.generate_content(prompt, instruction)

        return {"reply": reply_text}

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Chat engine encountered an error: {str(e)}")

@app.post("/api/generate-graph")
async def generate_graph(req: DocRequest):
    doc_data = db.get_document(req.doc_id)
    if not doc_data:
        raise HTTPException(status_code=404, detail="Document not found.")

    if doc_data["graph"]:
        return {
            "mermaid_code": doc_data["graph"]
        }

    try:
        prompt = f"Analyze this medical text and create a Mermaid.js mindmap showing core disease, symptoms, treatments. Start with 'mindmap' on line 1. Keep nodes short. Text: {doc_data['full_text'][:15000]}"
        response_text = ai_service.generate_content(prompt, "Output ONLY valid mermaid mindmap code. No markdown blocks.")
        
        mermaid_code = response_text.replace("```mermaid", "").replace("```", "").strip()
        db.update_document(req.doc_id, "graph", mermaid_code)

        return {
            "mermaid_code": mermaid_code
        }

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Graph Generation Failed")

@app.post("/api/generate-ppt")
async def generate_ppt(req: DocRequest):
    doc_data = db.get_document(req.doc_id)
    if not doc_data:
        raise HTTPException(status_code=404, detail="Document not found.")

    if doc_data["ppt"]:
        return {"slides": doc_data["ppt"]}

    try:
        prompt = f"Create a 5 slide presentation summarizing core findings. Schema: [{{ 'title': 'Title', 'bullets': ['pt1', 'pt2'], 'icon': 'ph-pill' }}]. Text: {doc_data['full_text'][:15000]}"
        response_text = ai_service.generate_content(prompt, "Output strictly a JSON array.", is_json=True)
        
        clean_text = response_text.replace("```json", "").replace("```", "").strip()
        slides = json.loads(clean_text)
        db.update_document(req.doc_id, "ppt", slides)
    
        return {
            "slides": slides
        }

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="PPT Generation Failed")

@app.post("/api/export-ppt")
async def export_ppt(req: DocRequest):
    doc_data = db.get_document(req.doc_id)
    if not doc_data or not doc_data.get("ppt"):
        raise HTTPException(status_code=404, detail="Please generate the PPT in the UI first.")
        
    try:
        slides_data = doc_data["ppt"]
        prs = Presentation()

        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1]) 
            title_shape = slide.shapes.title

            if title_shape:
                title_shape.text = slide_data.get("title", "Slide")

            if len(slide.shapes.placeholders) > 1:
                body_shape = slide.shapes.placeholders[1]
                t = body_shape.text_frame
                bullets = slide_data.get("bullets", [])

                if bullets:
                    t.clear()
                    for bullet in bullets:
                        p = t.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
            else:
                logger.warning(f"Layout for slide '{slide_data.get('title')}' is missing a content placeholder.")
                        
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)

        return StreamingResponse(
            ppt_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=MedAgent_Export.pptx"}
        )

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Export Failed")
